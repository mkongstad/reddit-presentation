import argparse

import os
import re
from dataclasses import dataclass
from io import BytesIO
from typing import Iterable, Tuple

import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from praw import Reddit


@dataclass
class SlideContent:
    image: BytesIO
    title: str


def generate_slide_content(subreddit_name: str, limit: int) -> Iterable[SlideContent]:
    reddit = Reddit(
        client_id=os.getenv("client_id"),
        client_secret=os.getenv("client_secret"),
        user_agent=os.getenv("user_agent"),
    )

    subreddit_name = reddit.subreddit(subreddit_name)

    def is_image(url: str) -> bool:
        return re.match("(.*\.(png|jpg))", url) is not None

    def download_image(url: str) -> Image:
        response = requests.get(url)
        return BytesIO(response.content)

    posts = subreddit_name.new(limit=limit)
    slide_content = []
    for post in posts:
        if not is_image(post.url):
            continue
        image = download_image(post.url)
        slide_content.append(SlideContent(image=image, title=post.title))

    return reversed(slide_content)


def calculate_placement(data: BytesIO) -> dict:
    placement = {"top": Inches(1.5), "left": 0}
    image = Image.open(data)
    w_max, h_max = Inches(10), Inches(6)
    w, h = map(lambda x: Inches(x), image.size)
    w_ratio, h_ratio = w_max / w, h_max / h

    if w_ratio < h_ratio:
        placement["width"] = w_max
    else:
        placement["height"] = h_max
    return placement


def create_presentation(
    title: str, description: str, slide_content: Iterable[SlideContent]
) -> None:
    presentation = Presentation()
    layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(layout)
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = description

    for content in slide_content:
        layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = content.title
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        placement = calculate_placement(content.image)
        slide.shapes.add_picture(content.image, **placement)
    presentation.save(f"{title}.pptx")


def parse_arguments() -> Tuple[str, int, str, str]:
    parser = argparse.ArgumentParser()
    parser.add_argument("--subreddit", type=str, required=True)
    parser.add_argument("--limit", type=int, required=True)
    parser.add_argument("--title", type=str, required=False)
    parser.add_argument("--description", type=str, required=False)
    args = parser.parse_args()
    title = args.title or args.subreddit
    description = args.description or ""
    return args.subreddit, args.limit, title, description


if __name__ == "__main__":
    subreddit, limit, title, description = parse_arguments()
    slide_content = generate_slide_content(subreddit, limit)
    create_presentation(title, description, slide_content)
